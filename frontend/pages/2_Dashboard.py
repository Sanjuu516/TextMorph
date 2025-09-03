import streamlit as st
import textstat
import plotly.express as px
import pandas as pd
import requests
# Imports for file parsing
import fitz  # PyMuPDF
import docx
from pptx import Presentation

# --- Check Login Status ---
# Ensures that only logged-in users can access this page
if not st.session_state.get('logged_in', False):
    st.error("You need to log in to view this page.")
    st.stop() # Halts the execution of the script

# --- Gemini API Call Function ---
def get_ai_analysis(text):
    """
    Calls the Gemini API to get a qualitative analysis of the provided text.
    Handles potential errors like missing API keys or network issues.
    """
    try:
        # Securely fetch the API key from Streamlit's secrets management
        api_key = st.secrets["GEMINI_API_KEY"]
    except FileNotFoundError:
        st.error("Please add your Gemini API key to .streamlit/secrets.toml")
        return None, None

    # Define the specific API endpoint URL
    api_url = f"https://generativelanguage.googleapis.com/v1beta/models/gemini-1.5-flash-latest:generateContent?key={api_key}"
    
    # Construct the prompt with clear instructions for the AI model
    prompt = f"""
    Analyze the following text for readability and tone. Provide two sections in your response:
    1.  **Analysis:** A brief, one-paragraph analysis of the text's complexity, style, and overall tone.
    2.  **Suggestions:** A bulleted list of 3-4 specific, actionable suggestions to improve the text's clarity and readability.

    Here is the text:
    ---
    {text}
    """

    # Structure the payload according to the Gemini API's requirements
    payload = {"contents": [{"parts": [{"text": prompt}]}]}
    headers = {'Content-Type': 'application/json'}
    
    try:
        # Make the POST request to the API
        response = requests.post(api_url, json=payload, headers=headers)
        
        # Check the HTTP status code from the response
        if response.status_code == 200:
            result = response.json()
            # Safely parse the nested JSON response to extract the generated text
            full_text = result.get('candidates', [{}])[0].get('content', {}).get('parts', [{}])[0].get('text', '')
            # Split the text into the two predefined sections for clean presentation
            if "Suggestions:" in full_text:
                analysis_part = full_text.split("Suggestions:")[0].replace("Analysis:", "").strip()
                suggestions_part = "Suggestions:" + full_text.split("Suggestions:")[1]
                return analysis_part, suggestions_part
            else:
                return full_text, "" # Return the full text if splitting fails
        # UPDATED: Specific handling for the "model overloaded" error
        elif response.status_code == 503:
            st.error("The AI model is currently overloaded. Please wait a moment and try again.", icon="⏳")
            return None, None
        else:
            # Handle other potential API errors
            st.error(f"Error calling Gemini API: {response.text}")
            return None, None
    except requests.exceptions.RequestException as e:
        # Handle network-level errors
        st.error(f"Could not connect to Gemini API: {e}")
        return None, None

# --- Helper function for soft classification ---
def calculate_level_percentages(grade):
    """
    Calculates a 'soft' classification percentage for each level based on the
    Flesch-Kincaid grade. This provides a more nuanced view than a hard classification.
    """
    # Define the ideal "peak" grade for each category
    b_peak, i_peak, a_peak = 5.0, 10.0, 15.0
    
    # Calculate the distance from the text's grade to each peak
    dist_b = abs(grade - b_peak)
    dist_i = abs(grade - i_peak)
    dist_a = abs(grade - a_peak)
    
    # Invert the distances so that closer scores get higher values
    # Add a small constant to prevent division by zero
    inv_dist_b = 1 / (dist_b + 0.1)
    inv_dist_i = 1 / (dist_i + 0.1)
    inv_dist_a = 1 / (dist_a + 0.1)
    
    # Normalize the values to get a percentage distribution
    total_inv_dist = inv_dist_b + inv_dist_i + inv_dist_a
    percentages = {
        "Beginner": (inv_dist_b / total_inv_dist) * 100,
        "Intermediate": (inv_dist_i / total_inv_dist) * 100,
        "Advanced": (inv_dist_a / total_inv_dist) * 100,
    }
    return percentages

# --- Main Dashboard UI ---
st.title("📝 Readability Analysis Dashboard")
st.markdown("Analyze your text for readability, complexity, and get AI-powered suggestions for improvement.")

# --- Input Methods ---
st.subheader("Provide Your Text")
# Use tabs for a clean UI to switch between input methods
input_method = st.tabs(["Paste Text", "Upload File"])

text_input = "" # Initialize an empty string to hold the user's text

# Tab 1: Pasting text directly
with input_method[0]:
    pasted_text = st.text_area("Paste your text here", height=250, key="pasted_text_area")
    if pasted_text:
        text_input = pasted_text

# Tab 2: Uploading a file
with input_method[1]:
    uploaded_file = st.file_uploader(
        "Upload a .txt, .pdf, .docx, or .pptx file",
        type=['txt', 'pdf', 'docx', 'pptx']
    )
    if uploaded_file is not None:
        try:
            # Check the MIME type to determine how to parse the file
            if uploaded_file.type == "text/plain":
                text_input = uploaded_file.getvalue().decode("utf-8")
            elif uploaded_file.type == "application/pdf":
                with fitz.open(stream=uploaded_file.read(), filetype="pdf") as doc:
                    text_input = "".join(page.get_text() for page in doc)
            elif uploaded_file.type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
                doc = docx.Document(uploaded_file)
                text_input = "\n".join([para.text for para in doc.paragraphs])
            elif uploaded_file.type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
                prs = Presentation(uploaded_file)
                text_input = ""
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text_input += shape.text + "\n"
            
            # Display the extracted text in a disabled text area for confirmation
            st.text_area("Extracted Content", text_input, height=250, disabled=True)
        except Exception as e:
            st.error(f"Error reading or processing file: {e}")

# --- Analysis Section ---
# This entire block only runs if text has been provided
if text_input:
    st.divider()
    st.subheader("Readability Scores")
    
    # Calculate the three main readability metrics using the textstat library
    fk_grade = textstat.flesch_kincaid_grade(text_input)
    gunning_fog = textstat.gunning_fog(text_input)
    smog_index = textstat.smog_index(text_input)

    # Use a three-column layout to display the scores neatly
    col1, col2, col3 = st.columns(3)
    col1.metric("Flesch-Kincaid Grade", f"{fk_grade:.2f}")
    col2.metric("Gunning Fog Index", f"{gunning_fog:.2f}")
    col3.metric("SMOG Index", f"{smog_index:.2f}")

    st.subheader("Text Composition Analysis")
    
    # Get the percentage breakdown from our helper function
    percentages = calculate_level_percentages(fk_grade)
    # Determine the primary level for a quick summary message
    primary_level = max(percentages, key=percentages.get)
    st.info(f"This text most closely resembles an **{primary_level}** level.")

    # Create a Pandas DataFrame, which is the required input format for Plotly
    chart_data = pd.DataFrame(list(percentages.items()), columns=["Level", "Percentage"])
    
    # Create the interactive bar chart using Plotly Express
    fig = px.bar(chart_data, 
                 x="Level", 
                 y="Percentage", 
                 text=chart_data['Percentage'].apply(lambda x: f'{x:.1f}%'), # Format text labels
                 color="Level",
                 color_discrete_map={  # Define custom colors for each category
                     "Beginner": "#5cb85c",
                     "Intermediate": "#f0ad4e",
                     "Advanced": "#d9534f"
                 })
    # Render the Plotly chart in the Streamlit app
    st.plotly_chart(fig, use_container_width=True)

    st.divider()
    st.subheader("🤖 AI-Powered Analysis & Suggestions")
    
    # The AI analysis is only performed when the user clicks this button
    if st.button("Get AI Analysis", use_container_width=True):
        # Display a spinner to provide visual feedback during the API call
        with st.spinner("Our AI is analyzing your text..."):
            analysis, suggestions = get_ai_analysis(text_input)
            if analysis and suggestions:
                st.success("Analysis Complete!")
                st.subheader("Analysis")
                st.write(analysis)
                st.subheader("Suggestions for Improvement")
                st.markdown(suggestions)
